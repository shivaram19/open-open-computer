"""Generate a Signal Network pitch deck for recruiting contributors."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).parent
OUTPUT = ROOT / "docs" / "Signal_Network_Pitch_Deck.pptx"
OUTPUT.parent.mkdir(parents=True, exist_ok=True)

# Brand colors
DARK = RGBColor(0x0F, 0x17, 0x2A)  # slate 900
ACCENT = RGBColor(0x06, 0xB6, 0xD4)  # cyan 500
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)  # slate 50
MUTED = RGBColor(0x94, 0xA3, 0xB8)  # slate 400


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(8.5), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = LIGHT
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.9), Inches(8.5), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    footer = slide.shapes.add_textbox(Inches(0.75), Inches(6.8), Inches(8.5), Inches(0.5))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Open-source · Indic-languages · GPU-powered"
    p.font.size = Pt(14)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.CENTER
    return slide


def add_bullet_slide(prs, title, bullets, subtitle=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT

    # Accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    y = Inches(0.6)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), y, Inches(9), Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = ACCENT
        p.font.bold = True
        y += Inches(0.45)

    title_box = slide.shapes.add_textbox(Inches(0.6), y, Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK
    y += Inches(1.0)

    content_box = slide.shapes.add_textbox(Inches(0.6), y, Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x33, 0x41, 0x55)
        p.space_after = Pt(14)
    return slide


def add_pipeline_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "How Signal Network Works"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = LIGHT

    steps = [
        ("1. Ingest", "yt-dlp + ffmpeg → clean audio"),
        ("2. ASR", "IndicConformer → word-level transcripts"),
        ("3. Translate", "IndicTrans2 → 22 Indic languages"),
        ("4. Clip", "Signal-aware viral-moment extraction"),
        ("5. Render", "Revideo / FFmpeg caption burn"),
        ("6. Publish", "Azure Blob + platform publishers"),
    ]

    x = Inches(0.6)
    y = Inches(1.7)
    w = Inches(2.7)
    h = Inches(1.1)
    gap_x = Inches(3.0)
    gap_y = Inches(1.5)

    for idx, (head, body) in enumerate(steps):
        col = idx % 3
        row = idx // 3
        left = x + col * gap_x
        top = y + row * gap_y

        box = slide.shapes.add_shape(1, left, top, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
        box.line.color.rgb = ACCENT
        box.line.width = Pt(1.5)

        tf = box.text_frame
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = head
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT

        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(13)
        p2.font.color.rgb = LIGHT

    arrow = slide.shapes.add_textbox(Inches(3.5), Inches(5.1), Inches(3), Inches(0.4))
    tf = arrow.text_frame
    p = tf.paragraphs[0]
    p.text = "One input → dozens of regional variants"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_table_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Built & Tested Stack"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK

    rows = 7
    cols = 3
    table = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.5), Inches(8.8), Inches(5.2)).table
    table.columns[0].width = Inches(2.4)
    table.columns[1].width = Inches(3.6)
    table.columns[2].width = Inches(2.8)

    headers = ["Layer", "Tool / Model", "Status"]
    data = [
        ["GPU VM", "Massed Compute A6000, CUDA 12.5", "Running"],
        ["ASR", "AI4Bharat IndicConformer 600M", "22-lang transcripts"],
        ["Translation", "AI4Bharat IndicTrans2 1B", "Indic→Indic"],
        ["Media Gen", "Bernini-R 1.3B Diffusers", "20/20 tests pass"],
        ["Pipeline", "Signal Network spine + fan-out", "35/35 tests pass"],
        ["Cloud Upload", "Azure Blob Storage", "Ready to wire"],
    ]

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = LIGHT
        p.font.bold = True
        p.font.size = Pt(16)

    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x33, 0x41, 0x55)
            if c == 2:
                p.font.color.rgb = RGBColor(0x05, 0x97, 0x55) if "pass" in val.lower() or "ready" in val.lower() else DARK
                p.font.bold = True
    return slide


def add_social_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Publish Everywhere"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = LIGHT

    platforms = [
        ("YouTube Shorts", "Global discovery + search"),
        ("Instagram Reels", "Visual-first regional shares"),
        ("TikTok", "Viral short-form momentum"),
        ("WhatsApp Status", "Hyper-local friend-to-friend spread"),
        ("Telegram Channels", "Community broadcast"),
        ("X / Twitter", "News + policy discourse"),
        ("LinkedIn", "Professional + NGO reach"),
        ("Facebook", "Older demographics + groups"),
    ]

    x = Inches(0.6)
    y = Inches(1.6)
    w = Inches(4.2)
    h = Inches(0.65)
    for idx, (name, desc) in enumerate(platforms):
        col = idx % 2
        row = idx // 2
        left = x + col * Inches(4.6)
        top = y + row * Inches(0.85)

        box = slide.shapes.add_shape(1, left, top, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
        box.line.color.rgb = ACCENT
        box.line.width = Pt(1)

        tf = box.text_frame
        tf.margin_left = Inches(0.1)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = MUTED

    return slide


def add_cta_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Join the Build"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = LIGHT
    p.alignment = PP_ALIGN.CENTER

    bullets = [
        "ML Engineers — optimize ASR/translation on GPU, ONNX, quantization",
        "Video / Frontend devs — real Revideo templates, caption styling, web dashboard",
        "Data / Linguists — improve Indic transcripts, region hashtags, demographic profiles",
        "DevOps / Cloud — Azure deployments, CI/CD, cost-efficient GPU scheduling",
        "Community — content testing, feedback loops, regional outreach",
    ]
    y = Inches(2.9)
    for text in bullets:
        box = slide.shapes.add_textbox(Inches(1.0), y, Inches(8), Inches(0.45))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"▸ {text}"
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT
        y += Inches(0.52)

    footer = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(8.5), Inches(0.8))
    tf = footer.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Repo: github.com/shivaram19/open-open-computer  ·  Signal Network folder"
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Share this deck on LinkedIn, X, Instagram, WhatsApp, Telegram — anywhere builders gather."
    p2.font.size = Pt(12)
    p2.font.color.rgb = MUTED
    p2.alignment = PP_ALIGN.CENTER
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Signal Network",
        "Open-source hyper-local edutainment at scale\nOne video → N languages × M formats × P platforms",
    )

    add_bullet_slide(
        prs,
        "The Problem",
        [
            "India has 22 scheduled languages and thousands of dialects.",
            "Civic and economic information rarely reaches people in their native tongue.",
            "Manual translation, dubbing, and short-form editing don’t scale.",
            "Creators and NGOs need a deterministic, reusable pipeline.",
        ],
        subtitle="WHY THIS MATTERS",
    )

    add_bullet_slide(
        prs,
        "The Solution",
        [
            "Drop one source video — get dozens of regional short-form clips.",
            "ASR, translation, clip extraction, caption burn, and publishing are automated.",
            "Demographic engine matches topics to regions, languages, and best posting times.",
            "Fully open-source: models, code, and templates are transparent and forkable.",
        ],
        subtitle="WHAT WE’RE BUILDING",
    )

    add_pipeline_slide(prs)

    add_table_slide(prs)

    add_social_slide(prs)

    add_bullet_slide(
        prs,
        "What’s Already Working",
        [
            "GPU VM with CUDA 12.5 + cuDNN 9 deployed and validated.",
            "Bernini image generation: 20/20 tests passing.",
            "IndicConformer ASR + IndicTrans2 translation on GPU.",
            "Signal Network spine, clip extraction, fan-out: 35/35 tests passing.",
            "Azure Blob Storage publisher implemented and ready for credentials.",
        ],
        subtitle="CURRENT STATUS",
    )

    add_bullet_slide(
        prs,
        "Immediate Next Steps",
        [
            "Wire real Azure credentials and upload the first batch of clips.",
            "Replace FFmpeg subtitle burn with true Revideo 9:16 templates.",
            "Add platform publishers: YouTube Shorts, Instagram Reels, TikTok.",
            "Build a lightweight web dashboard for upload + monitoring.",
            "Crowdsource demographic profiles and regional signal data sources.",
        ],
        subtitle="ROADMAP",
    )

    add_cta_slide(prs)

    prs.save(OUTPUT)
    print(f"Saved pitch deck to {OUTPUT}")


if __name__ == "__main__":
    main()
